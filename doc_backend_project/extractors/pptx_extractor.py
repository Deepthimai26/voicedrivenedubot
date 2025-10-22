from pptx import Presentation
from PIL import Image
from io import BytesIO
import os

def extract_from_pptx(file_path, output_dir="extracted_data/pptx"):
    os.makedirs(output_dir, exist_ok=True)
    prs = Presentation(file_path)
    text_content = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            # ---- Extract Text ----
            if hasattr(shape, "text"):
                text_content.append(shape.text)

            # ---- Extract Images ----
            if shape.shape_type == 13:  # Picture
                image = shape.image
                image_bytes = image.blob
                img = Image.open(BytesIO(image_bytes))
                img_filename = f"slide{slide_num}_img{len(os.listdir(output_dir)) + 1}.png"
                img.save(os.path.join(output_dir, img_filename))
            if shape.has_table:
                table = shape.table
                table_data = [[cell.text.strip() for cell in row.cells] for row in table.rows]
    # Optionally save table_data or return along with text/images


    return {"text": "\n".join(text_content), "image_dir": output_dir}
