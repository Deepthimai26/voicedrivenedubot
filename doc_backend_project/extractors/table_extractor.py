import os
import pandas as pd
import pdfplumber
import tabula
from docx import Document
from pptx import Presentation

def extract_tables_from_pdf(file_path: str):
    """Extract tables from PDF using Camelot or pdfplumber"""
    tables = []

    # Try with Camelot first
    try:
        dfs = tabula.read_pdf(file_path, pages="all", multiple_tables=True)
        if dfs:  # check if tabula returned any tables
            tables.extend(dfs) 
    except Exception:
        # Fallback to pdfplumber
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                for table in page.extract_tables():
                    df = pd.DataFrame(table[1:], columns=table[0])
                    tables.append(df)
    return tables


def extract_tables_from_docx(file_path: str):
    """Extract tables from Word file"""
    doc = Document(file_path)
    tables = []
    for t in doc.tables:
        df = pd.DataFrame([[cell.text for cell in row.cells] for row in t.rows])
        tables.append(df)
    return tables


def extract_tables_from_pptx(file_path: str):
    """Extract tables from PowerPoint file"""
    prs = Presentation(file_path)
    tables = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                df = pd.DataFrame(
                    [
                        [cell.text.strip() for cell in row.cells]
                        for row in table.rows
                    ]
                )
                tables.append(df)
    return tables


def save_tables_to_csv(tables, output_dir, base_name):
    """Save extracted tables to CSV"""
    os.makedirs(output_dir, exist_ok=True)
    saved_files = []
    for i, df in enumerate(tables):
        file_name = os.path.join(output_dir, f"{base_name}_table_{i+1}.csv")
        df.to_csv(file_name, index=False)
        saved_files.append(file_name)
    return saved_files
